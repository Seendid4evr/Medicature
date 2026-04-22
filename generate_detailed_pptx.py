import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()
# Set presentation size to 16:9
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Layouts
blank_slide_layout = prs.slide_layouts[6]
title_slide_layout = prs.slide_layouts[0]
title_content_layout = prs.slide_layouts[1]

# 1. Add Existing Slides (Screenshots from PopAI)
image_folder = r"Z:\Xampp\htdocs\medicure\sss"
images = [f for f in os.listdir(image_folder) if f.startswith("Screenshot") and f.endswith(".png")]
images.sort()

for img_file in images:
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = os.path.join(image_folder, img_file)
    # Add image to fit the whole slide
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

# 2. Add the WBS Mermaid Diagram Slide
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title
title.text = "Work Breakdown Structure (WBS) - Overview"
img_path = os.path.join(image_folder, "mermaid-diagram-2026-03-14-105219.png")
# Add image scaled
slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(14))

# 3. Add Detailed WBS Slides
wbs_details = [
    ("Phase 1: Initiation", [
        "Objective: Define project scope and feasibility for a zero-friction web solution.",
        "1.1 Requirement Gathering: Analyze medication non-adherence issues, interview stakeholders, establish core requirement of eliminating app-store friction.",
        "1.2 Feasibility Study: Research web browser capabilities to handle background audio alarms via service workers.",
        "1.3 Project Charter: Formalize project scope, rejecting directory bloat (unlike competitors), and securing stakeholder approval."
    ]),
    ("Phase 2: Design", [
        "Objective: Architect the resilient, mobile-first technical foundation.",
        "2.1 Architecture Planning: Design Entity-Relationship Diagrams (ERDs) mapping Users -> Medicines -> Schedules -> Reminders.",
        "2.2 UI/UX Wireframing: Create low-fidelity mockups emphasizing large, accessible typography and 'hit areas' for elderly users.",
        "2.3 Interactive Prototyping: Build high-fidelity Figma screens with a mobile-first responsive design strategy."
    ]),
    ("Phase 3: Development", [
        "Objective: Build the core application and proprietary alarm logic.",
        "3.1 Frontend Dashboard Development: Code the responsive UI using HTML/CSS/JS for cross-platform browser support.",
        "3.2 Backend CRUD APIs: Develop secure PHP/SQL endpoints to handle sensitive medical data creation, reading, updating, and deletion.",
        "3.3 The Notification Engine (Core Intellectual Property): Engineer JavaScript Service Workers & Web Audio API integration to trigger alarms reliably."
    ]),
    ("Phase 4: Implementation", [
        "Objective: Rigorously test and deploy the stable product.",
        "4.1 Integration Testing: Connect Frontend UI with Backend schedules, ensuring strict timestamp accuracy.",
        "4.2 User Acceptance Testing (UAT): Simulate hundreds of medication drops across Chrome, Safari, and Firefox browsers to validate alarm triggers.",
        "4.3 Live Deployment: Secure domain registration, HTTPS certificate installation, and push code to live web hosting servers."
    ]),
    ("Phase 5: Closure", [
        "Objective: Finalize documentation and prepare for future scaling.",
        "5.1 System Documentation: Export final ERDs, Mermaid WBS diagrams, and Use Case architectures to project repositories.",
        "5.2 Project Handover: Establish a scalable framework and knowledge base for onboarding future developers."
    ])
]

for title_text, bullet_points in wbs_details:
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = title_text
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(28)
        else:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.level = 1 if i > 0 else 0

output_path = r"Z:\Xampp\htdocs\medicure\docs\Medicature_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
